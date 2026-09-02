# -*- coding: utf-8 -*-
"""日本語ビジネスデッキ用の python-pptx ヘルパ（16:9 / 13.333 x 7.5 インチ）。

使い方
------
このファイルをスクラッチパッドにコピーし、build.py から import する。

    from deck_lib import *

    prs = new_deck()

    s = slide_new(prs, dark=True)                       # 表紙（ダーク地）
    box(s, ML, 2.4, 11.0, 1.5, "見出し", size=40, bold=True, color=WHITE, spacing=1.22)
    notes(s, "ここで話すこと")

    s = slide_new(prs)                                  # 本編（ライト地）
    head(s, "1文の主張をここに書く", sub="補足", badge="要検証", badge_fill=AMBER)
    card(s, ML, 2.2, 5.35, 2.6)                         # 角丸カード
    box(s, ML + 0.4, 2.5, 4.6, 1.8, "本文", size=13)
    table(s, ML, 5.0, CW, [["列A", "列B"], ["1", "2"]], [6.0, 5.833])
    footnote(s, "※ 出典・但し書き")
    pagenum(s, 2)

    prs.save("deck.pptx")

box() のテキスト引数は3通り
    "ただの文字列"
    [("太字", {"bold": True}), ("続き", {"color": MUTED})]      # 1段落・複数run
    [[("1段落目", {})], [("2段落目", {"size": 12})]]            # 複数段落

色は PowerPoint 側で解釈されるので、ダーク地には *_D / 明るい版を使う。
数値の根拠と落とし穴は references/layout.md を見ること。
"""
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette
# 内容に合わせて差し替える前提。意味と色を結びつけると通しのモチーフになる。
INK      = RGBColor(0x0F, 0x2C, 0x3A)   # ダーク地
INK_CARD = RGBColor(0x1B, 0x44, 0x56)   # ダーク地の上のカード
PAPER    = RGBColor(0xFF, 0xFF, 0xFF)
TINT     = RGBColor(0xF0, 0xF4, 0xF6)   # ライト地のカード
TINT2    = RGBColor(0xE3, 0xEA, 0xED)   # 表のヘッダ行
TEXT     = RGBColor(0x18, 0x22, 0x28)
MUTED    = RGBColor(0x5E, 0x72, 0x7B)
MUTED_D  = RGBColor(0x9D, 0xB4, 0xBE)   # ダーク地のミュート
GREEN    = RGBColor(0x14, 0x6E, 0x4C);  GREEN_T = RGBColor(0xE4, 0xF1, 0xEB)
RED      = RGBColor(0xA0, 0x33, 0x28);  RED_T   = RGBColor(0xF8, 0xE9, 0xE7)
AMBER    = RGBColor(0xA8, 0x71, 0x0F);  AMBER_T = RGBColor(0xFB, 0xF1, 0xDF)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
# ダーク地に置くときはこちら（濃い色をダーク地に置くと読めない）
GREEN_L  = RGBColor(0x5F, 0xC7, 0x9B)
RED_L    = RGBColor(0xE8, 0x8B, 0x7E)
AMBER_L  = RGBColor(0xE3, 0xAE, 0x4A)

FONT = "Yu Gothic"      # Windows 8.1+ / Mac Office の両方にある

W, H = 13.333, 7.5
ML, MR = 0.75, 0.75
CW = W - ML - MR        # 本文幅 11.833


# ---------------------------------------------------------------- internals
def _typeface(run, name):
    """latin だけでなく ea（東アジア）にも書かないと日本語のフォントが変わらない。"""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def noshadow(shape):
    shape.shadow.inherit = False


# ---------------------------------------------------------------- API
def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = In(W), In(H)
    return prs


def slide_new(prs, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])   # 白紙レイアウト
    if dark:
        bg = s.background.fill
        bg.solid()
        bg.fore_color.rgb = INK
    return s


def box(slide, x, y, w, h, text="", size=14, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0,
        space_after=0, wrap=True):
    """テキストボックス。内側余白は0にしてある（図形と左端を揃えるため）。"""
    tb = slide.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    if isinstance(text, str):
        paras = [[(text, {})]]
    elif text and isinstance(text[0], tuple):
        paras = [text]
    else:
        paras = text

    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if space_after:
            p.space_after = Pt(space_after)
        for s, o in runs:
            r = p.add_run()
            r.text = s
            r.font.size = Pt(o.get("size", size))
            r.font.bold = o.get("bold", bold)
            r.font.color.rgb = o.get("color", color)
            _typeface(r, FONT)
    return tb


def rect(slide, x, y, w, h, fill=None, line=None, lw=1.0,
         shape=MSO_SHAPE.RECTANGLE, radius=None):
    sh = slide.shapes.add_shape(shape, In(x), In(y), In(w), In(h))
    noshadow(sh)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.text_frame.word_wrap = True
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = radius
    return sh


def card(slide, x, y, w, h, fill=TINT, line=None):
    """角丸カード。3枚並べるなら w=3.72 / ピッチ3.98。"""
    return rect(slide, x, y, w, h, fill=fill, line=line,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.045)


def pill(slide, x, y, label, fill, w=None, tcol=WHITE, size=11):
    """状態バッジ。通しのモチーフとして使うと内容と意味が結びつく。"""
    if w is None:
        w = 0.30 + len(label) * (size / 72.0) * 1.02
    sh = rect(slide, x, y, w, 0.34, fill=fill,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = tcol
    _typeface(r, FONT)
    return sh


def head(slide, text, dark=False, sub=None, badge=None, badge_fill=None):
    """見出し＝1文の主張（アサーション・エビデンス型）。26pt・2行までを見込む。"""
    box(slide, ML, 0.52, CW - (2.6 if badge else 0), 1.05, text,
        size=26, bold=True, color=WHITE if dark else TEXT, spacing=1.18)
    if sub:
        box(slide, ML, 1.62, CW, 0.34, sub, size=13,
            color=MUTED_D if dark else MUTED, spacing=1.1)
    if badge:
        pill(slide, W - MR - 2.35, 0.60, badge, badge_fill, w=2.35)


def table(slide, x, y, w, rows, colw, head_fill=TINT2, row_h=0.42,
          size=12, head_size=12, cell_colors=None):
    """rows[0] がヘッダ行。cell_colors[(r, c)] = (文字色, セル地色)。
       高さは row_h × 行数。この値で脚注の y を決めること。"""
    nr, nc = len(rows), len(rows[0])
    gf = slide.shapes.add_table(nr, nc, In(x), In(y), In(w), In(row_h * nr))
    tbl = gf.table
    tblPr = tbl._tbl.find(qn("a:tblPr"))
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")
    for e in tblPr.findall(qn("a:tableStyleId")):
        tblPr.remove(e)
    sid = tblPr.makeelement(qn("a:tableStyleId"), {})
    sid.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"   # No Style, No Grid
    tblPr.append(sid)

    for i, cw in enumerate(colw):
        tbl.columns[i].width = In(cw)
    for r in range(nr):
        tbl.rows[r].height = In(row_h)

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left, cell.margin_right = In(0.12), In(0.10)
            cell.margin_top, cell.margin_bottom = In(0.05), In(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            fcol, bcol = (TEXT, head_fill if r == 0 else None)
            if cell_colors and (r, c) in cell_colors:
                fcol, bcol = cell_colors[(r, c)]
            cell.fill.solid()
            cell.fill.fore_color.rgb = bcol if bcol is not None else (
                PAPER if r % 2 else TINT)
            p = cell.text_frame.paragraphs[0]
            cell.text_frame.word_wrap = True
            p.alignment = PP_ALIGN.RIGHT if (c > 0 and val.replace(
                ",", "").replace("%", "").isdigit()) else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = val
            run.font.size = Pt(head_size if r == 0 else size)
            run.font.bold = (r == 0)
            run.font.color.rgb = fcol
            _typeface(run, FONT)
    return tbl


def footnote(slide, text, dark=False, y=6.92):
    """幅を本文幅−0.95 にしてあるのは、右下のページ番号と重ねないため。"""
    box(slide, ML, y, CW - 0.95, 0.4, text, size=10,
        color=MUTED_D if dark else MUTED, spacing=1.1)


def pagenum(slide, n, dark=False):
    box(slide, W - MR - 0.7, 6.92, 0.7, 0.3, str(n), size=10,
        color=MUTED_D if dark else MUTED, align=PP_ALIGN.RIGHT)


def notes(slide, text):
    """話しながら見せる資料では、説明はスライド本文ではなくここに書く。"""
    slide.notes_slide.notes_text_frame.text = text
