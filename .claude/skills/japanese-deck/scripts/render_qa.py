# -*- coding: utf-8 -*-
"""pptx を PIL で近似レンダリングし、はみ出しと枠外を機械検出する。

LibreOffice も node も無い環境で目視QAするための道具。
PowerPoint 本体の AppleScript 経由の PDF 書き出しは返ってこないことがあるので、
1回試して駄目ならこちらに切り替える。

    python3 render_qa.py deck.pptx [出力ディレクトリ]

png/slide-01.png ... を書き出し、警告を標準出力に一覧する。
**警告ゼロは合格ではない。** 出た画像を必ず全部見ること。

限界（伝えるべきこと）
  - フォントは環境にあるCJKフォントで代用する。字面の印象は実物と違う
  - 太字はストロークで擬似的に描く
  - 折り返しは自前計算。PowerPoint の禁則処理とは完全には一致しない
  → 行数とはみ出しの判定は実用上一致するが、最終的な見え方の保証ではない
"""
import sys, re, os, glob
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.oxml.ns import qn

DPI = 110

FONT_CANDIDATES = [
    "/Library/Fonts/Arial Unicode.ttf",
    "/System/Library/Fonts/Hiragino Sans GB.ttc",
    "/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc",
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    "/usr/share/fonts/truetype/fonts-japanese-gothic.ttf",
    "C:/Windows/Fonts/YuGothM.ttc",
    "C:/Windows/Fonts/meiryo.ttc",
]


def find_font():
    for p in FONT_CANDIDATES:
        if os.path.exists(p):
            return p
    for pat in ("/System/Library/Fonts/**/*.ttc", "/usr/share/fonts/**/*.tt[cf]"):
        for p in glob.glob(pat, recursive=True):
            if re.search(r"(hiragino|noto.*cjk|gothic|yugoth)", p, re.I):
                return p
    raise SystemExit("CJKフォントが見つからない。FONT_CANDIDATES にパスを足すこと。")


FONT_PATH = find_font()
_fc = {}


def F(pt):
    px = max(6, int(round(pt * DPI / 72.0)))
    if px not in _fc:
        _fc[px] = ImageFont.truetype(FONT_PATH, px)
    return _fc[px]


def E(v):
    return v / 914400.0 * DPI


def rgb(c):
    try:
        if c and c.type is not None and c.rgb is not None:
            return tuple(c.rgb)
    except Exception:
        pass
    return None


def shape_fill(sh):
    try:
        f = sh.fill
        if f.type is not None and int(f.type) == 1:
            return rgb(f.fore_color)
    except Exception:
        pass
    return None


def shape_line(sh):
    try:
        if sh.line.fill.type is not None and int(sh.line.fill.type) == 1:
            return rgb(sh.line.color)
    except Exception:
        pass
    return None


ASCIIRUN = re.compile(r"[0-9A-Za-z\.\,\%\$\€\'\-/\(\)]")


def wrap(runs, maxw):
    """CJKは任意位置、ラテンは語単位で折り返す近似。"""
    lines, cur, curw = [], [], 0.0
    for text, size, bold, color in runs:
        f = F(size)
        i = 0
        while i < len(text):
            ch = text[i]
            if ch == "\n":
                lines.append(cur); cur, curw = [], 0.0; i += 1; continue
            if ASCIIRUN.match(ch):
                j = i
                while j < len(text) and (ASCIIRUN.match(text[j]) or text[j] == " "):
                    j += 1
                tok = text[i:j]
            else:
                tok, j = ch, i + 1
            w = f.getlength(tok)
            if curw + w > maxw and cur:
                lines.append(cur); cur, curw = [], 0.0
                if tok.startswith(" "):
                    tok = tok.lstrip(" "); w = f.getlength(tok)
                if not tok:
                    i = j; continue
            if cur and cur[-1][1:] == (size, bold, color):
                cur[-1] = (cur[-1][0] + tok, size, bold, color)
            else:
                cur.append((tok, size, bold, color))
            curw += w
            i = j
    lines.append(cur)
    return lines


def draw_tf(d, tf, x, y, w, h, warn, label):
    ml, mr = E(tf.margin_left or 0), E(tf.margin_right or 0)
    mt, mb = E(tf.margin_top or 0), E(tf.margin_bottom or 0)
    tx, ty = x + ml, y + mt
    tw = max(4, w - ml - mr)
    wrapon = tf.word_wrap is not False

    blocks = []
    for p in tf.paragraphs:
        runs = [(r.text, r.font.size.pt if r.font.size else 18,
                 bool(r.font.bold), rgb(r.font.color) or (0, 0, 0)) for r in p.runs]
        if not runs:
            runs = [("", 12, False, (0, 0, 0))]
        blocks.append((wrap(runs, tw if wrapon else 10 ** 6),
                       p.line_spacing or 1.0,
                       p.space_after.pt if p.space_after else 0,
                       p.alignment))

    total = 0.0
    for lines, ls, sa, _ in blocks:
        for ln in lines:
            total += max([s for _, s, _, _ in ln] or [12]) * DPI / 72.0 * ls
        total += sa * DPI / 72.0

    avail = h - mt - mb
    anchor = tf.vertical_anchor
    cy = ty
    if anchor is not None and int(anchor) == 3:
        cy = y + (h - total) / 2
    elif anchor is not None and int(anchor) == 4:
        cy = y + h - mb - total
    if total > avail + 1.5:
        warn.append(f"{label}: はみ出し {total - avail:.0f}px (要{total:.0f} / 枠{avail:.0f})")

    for lines, ls, sa, al in blocks:
        for ln in lines:
            mx = max([s for _, s, _, _ in ln] or [12])
            lh = mx * DPI / 72.0 * ls
            lw = sum(F(s).getlength(t) for t, s, _, _ in ln)
            sx = tx
            if al is not None and int(al) == 2:
                sx = tx + (tw - lw) / 2
            elif al is not None and int(al) == 3:
                sx = tx + tw - lw
            base = cy + (lh - mx * DPI / 72.0) / 2
            for t, s, b, c in ln:
                d.text((sx, base), t, font=F(s), fill=c,
                       stroke_width=1 if (b and s >= 13) else 0, stroke_fill=c)
                sx += F(s).getlength(t)
            cy += lh
        cy += sa * DPI / 72.0


def main():
    src = sys.argv[1] if len(sys.argv) > 1 else sys.exit(__doc__)
    out = sys.argv[2] if len(sys.argv) > 2 else "png"
    os.makedirs(out, exist_ok=True)

    prs = Presentation(src)
    SW, SH = int(E(prs.slide_width)), int(E(prs.slide_height))
    allwarn = []

    for idx, slide in enumerate(prs.slides, 1):
        bg = (255, 255, 255)
        try:
            b = slide.background.fill
            if b.type is not None and int(b.type) == 1:
                bg = rgb(b.fore_color) or bg
        except Exception:
            pass
        img = Image.new("RGB", (SW, SH), bg)
        d = ImageDraw.Draw(img)
        warn = []

        for sh in slide.shapes:
            x, y, w, h = E(sh.left), E(sh.top), E(sh.width), E(sh.height)
            if sh.has_table:
                cy = y
                colw = [E(c.width) for c in sh.table.columns]
                for r, row in enumerate(sh.table.rows):
                    rh, cx = E(row.height), x
                    for c, cell in enumerate(row.cells):
                        fc = None
                        try:
                            if cell.fill.type is not None and int(cell.fill.type) == 1:
                                fc = rgb(cell.fill.fore_color)
                        except Exception:
                            pass
                        if fc:
                            d.rectangle([cx, cy, cx + colw[c], cy + rh], fill=fc)
                        draw_tf(d, cell.text_frame, cx, cy, colw[c], rh, warn, f"表 r{r}c{c}")
                        cx += colw[c]
                    cy += rh
                continue

            f, l = shape_fill(sh), shape_line(sh)
            if f or l:
                prst = ""
                try:
                    prst = sh._element.find(".//" + qn("a:prstGeom")).get("prst")
                except Exception:
                    pass
                if prst == "roundRect":
                    d.rounded_rectangle([x, y, x + w, y + h],
                                        radius=min(h, w) * 0.10, fill=f, outline=l)
                elif prst == "rightArrow":
                    d.polygon([(x, y + h * .3), (x + w * .6, y + h * .3), (x + w * .6, y),
                               (x + w, y + h / 2), (x + w * .6, y + h),
                               (x + w * .6, y + h * .7), (x, y + h * .7)], fill=f)
                else:
                    d.rectangle([x, y, x + w, y + h], fill=f, outline=l)

            if sh.has_text_frame and sh.text_frame.text.strip():
                draw_tf(d, sh.text_frame, x, y, w, h, warn,
                        f"「{sh.text_frame.text[:16]}」")
            if x < 0 or y < 0 or x + w > SW + 1 or y + h > SH + 1:
                lbl = sh.text_frame.text[:14] if sh.has_text_frame else "図形"
                warn.append(f"枠外: 「{lbl}」")

        img.save(f"{out}/slide-{idx:02d}.png")
        if warn:
            allwarn.append((idx, warn))

    print(f"rendered {len(prs.slides._sldIdLst)} slides -> {out}/  (font: {os.path.basename(FONT_PATH)})")
    if allwarn:
        print("\n=== 警告 ===")
        for i, ws in allwarn:
            print(f"\n[slide {i}]")
            for w in ws:
                print("  -", w)
    else:
        print("警告なし（※ これは合格ではない。画像を全部見ること）")


if __name__ == "__main__":
    main()
